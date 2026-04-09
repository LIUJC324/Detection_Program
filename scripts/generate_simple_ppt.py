from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
ARCH_IMAGE = PROJECT_ROOT / "docs" / "project_architecture.png"
OUTPUT_PPT = PROJECT_ROOT / "docs" / "rgbt_uav_detection_简易汇报.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLOR_BG = RGBColor(247, 249, 252)
COLOR_TEXT = RGBColor(25, 42, 86)
COLOR_SUB = RGBColor(72, 101, 129)
COLOR_LINE = RGBColor(216, 225, 232)
COLOR_BLUE = RGBColor(76, 154, 255)
COLOR_GREEN = RGBColor(99, 182, 111)
COLOR_ORANGE = RGBColor(245, 166, 35)
COLOR_RED = RGBColor(233, 87, 63)
COLOR_TEAL = RGBColor(39, 174, 162)
COLOR_PURPLE = RGBColor(127, 90, 240)
COLOR_CARD = RGBColor(255, 255, 255)


def add_full_background(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(SLIDE_W),
        Inches(SLIDE_H),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_top_bar(slide, color: RGBColor) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(SLIDE_W),
        Inches(0.18),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int,
    color: RGBColor = COLOR_TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
    valign=MSO_ANCHOR.TOP,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    accent: RGBColor,
) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_LINE
    card.line.width = Pt(1.2)
    card.shadow.inherit = False

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.14),
        Inches(height),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    add_textbox(slide, left + 0.28, top + 0.18, width - 0.42, 0.42, title, font_size=20, bold=True)
    add_textbox(
        slide,
        left + 0.28,
        top + 0.70,
        width - 0.42,
        height - 0.84,
        body,
        font_size=12,
        color=COLOR_SUB,
    )


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, COLOR_BG)
    add_top_bar(slide, COLOR_BLUE)

    hero = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(5.2),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.line.color.rgb = COLOR_LINE

    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(1.15),
        Inches(2.15),
        Inches(0.48),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(230, 241, 255)
    tag.line.fill.background()
    add_textbox(
        slide,
        1.02,
        1.22,
        2.0,
        0.3,
        "项目简易汇报",
        font_size=16,
        color=COLOR_BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        0.95,
        1.9,
        10.8,
        0.9,
        "RGBT UAV Detection",
        font_size=28,
        bold=True,
    )
    add_textbox(
        slide,
        0.95,
        2.65,
        10.8,
        0.7,
        "基于 RGB 与 Thermal 双模态输入的无人机视角小目标检测方案",
        font_size=20,
        color=COLOR_SUB,
    )

    add_card(
        slide,
        0.95,
        3.55,
        3.6,
        1.45,
        "汇报重点",
        "展示初级阶段架构图，并聚焦创新点、关键难题与解决路径。",
        COLOR_BLUE,
    )
    add_card(
        slide,
        4.75,
        3.55,
        3.6,
        1.45,
        "项目定位",
        "先打通数据、训练、推理、服务闭环，再围绕多模态与小目标性能做优化。",
        COLOR_GREEN,
    )
    add_card(
        slide,
        8.55,
        3.55,
        3.1,
        1.45,
        "适用场景",
        "课程汇报、阶段答辩、技术方案说明。",
        COLOR_ORANGE,
    )

    add_textbox(
        slide,
        0.95,
        5.45,
        11.0,
        0.45,
        "内容结构：项目定位 -> 整体架构 -> 创新点 -> 难题与理论进度",
        font_size=13,
        color=COLOR_SUB,
    )


def build_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, COLOR_BG)
    add_top_bar(slide, COLOR_GREEN)
    add_textbox(slide, 0.6, 0.35, 6.4, 0.45, "项目整体架构（初级阶段）", font_size=24, bold=True)
    add_textbox(
        slide,
        0.6,
        0.72,
        8.8,
        0.32,
        "架构图直接对应仓库中的数据层、模型层、训练评估层、推理导出层和服务集成层。",
        font_size=12,
        color=COLOR_SUB,
    )

    slide.shapes.add_picture(str(ARCH_IMAGE), Inches(0.45), Inches(1.1), width=Inches(8.85))

    add_card(
        slide,
        9.6,
        1.2,
        3.05,
        2.2,
        "当前阶段判断",
        "架构设计与代码主链路已经具备，适合做方案展示与后续实验规划；正式训练权重与指标结果仍待补齐。",
        COLOR_GREEN,
    )
    add_card(
        slide,
        9.6,
        3.6,
        3.05,
        1.75,
        "当前汇报口径",
        "更适合表述为“技术方案已成型、工程骨架已搭好”，不要直接表述成“性能验证已经完成”。",
        COLOR_RED,
    )
    add_card(
        slide,
        9.6,
        5.55,
        3.05,
        1.1,
        "下一步",
        "真实数据接入 -> 基线训练 -> 指标验证 -> 部署联调",
        COLOR_BLUE,
    )


def build_innovations(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, COLOR_BG)
    add_top_bar(slide, COLOR_ORANGE)
    add_textbox(slide, 0.6, 0.35, 6.0, 0.45, "方案创新点", font_size=24, bold=True)
    add_textbox(
        slide,
        0.6,
        0.72,
        9.5,
        0.32,
        "创新重点不是单个模块堆叠，而是围绕多模态互补和无人机小目标场景形成一条完整检测链路。",
        font_size=12,
        color=COLOR_SUB,
    )

    add_card(
        slide,
        0.7,
        1.25,
        3.9,
        2.05,
        "1. 双流轻量骨干",
        "RGB 与 Thermal 分支先独立提特征，再进入中层融合，避免过早混合带来的模态干扰；轻量卷积结构兼顾实时性。",
        COLOR_BLUE,
    )
    add_card(
        slide,
        4.75,
        1.25,
        3.9,
        2.05,
        "2. 跨模态注意力融合",
        "融合时同时使用 RGB 特征、Thermal 特征和二者差异信息，动态学习“当前位置更该相信哪个模态”。",
        COLOR_PURPLE,
    )
    add_card(
        slide,
        8.8,
        1.25,
        3.8,
        2.05,
        "3. 小目标增强链路",
        "BiFPN 负责多尺度语义传递，小目标精炼头补足局部细节，FCOS 则减少 anchor 设计对小目标的限制。",
        COLOR_ORANGE,
    )
    add_card(
        slide,
        0.7,
        3.65,
        5.95,
        1.65,
        "4. 面向落地的工程闭环",
        "训练、评估、导出、服务接口都在同一仓库中，便于后续从实验代码过渡到前后端联调与部署。",
        COLOR_TEAL,
    )

    add_card(
        slide,
        6.9,
        3.65,
        5.7,
        1.65,
        "一句话概括",
        "方案不是只追求更复杂的网络，而是让“全天候多模态检测 + 小目标增强 + 可部署”同时成立。",
        COLOR_GREEN,
    )


def build_challenges_and_progress(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide, COLOR_BG)
    add_top_bar(slide, COLOR_TEAL)
    add_textbox(slide, 0.6, 0.35, 8.5, 0.45, "关键难题、解决方法与理论进度", font_size=24, bold=True)
    add_textbox(
        slide,
        0.6,
        0.72,
        8.8,
        0.32,
        "左侧回答“为什么难、怎么解”，右侧说明项目当前更接近哪一个推进阶段。",
        font_size=12,
        color=COLOR_SUB,
    )

    add_card(
        slide,
        0.7,
        1.2,
        7.45,
        1.45,
        "难题 1：无人机视角下目标太小，卷积下采样后细节容易丢失",
        "解决方法：保留多尺度特征，用 BiFPN 做高低层融合，再用小目标精炼头补充局部边缘和上下文信息。",
        COLOR_ORANGE,
    )
    add_card(
        slide,
        0.7,
        2.9,
        7.45,
        1.45,
        "难题 2：夜间、低照、逆光等条件下，单一 RGB 模态不稳定",
        "解决方法：采用 RGB-T 双流输入，并通过跨模态注意力实现动态加权，在不同环境下自动强调更可靠的模态。",
        COLOR_PURPLE,
    )
    add_card(
        slide,
        0.7,
        4.6,
        7.45,
        1.45,
        "难题 3：密集遮挡和尺度变化大，容易漏检、误检",
        "解决方法：使用 Anchor-free 的 FCOS 降低先验框设计负担，同时结合多尺度特征增强提升密集场景定位稳定性。",
        COLOR_RED,
    )

    add_card(
        slide,
        8.45,
        1.2,
        4.2,
        4.85,
        "理论推进进度",
        "[已完成] 架构设计与模块拆分\n"
        "[已完成] 训练、评估、导出、服务代码骨架\n"
        "[进行中] 数据接入与标注质检\n"
        "[待推进] 基线训练与 best.pt 产出\n"
        "[待推进] 指标优化、模型导出与系统联调",
        COLOR_TEAL,
    )

    add_textbox(
        slide,
        8.65,
        5.55,
        3.8,
        0.52,
        "结论：当前最适合把它讲成“技术方案已成型，实验验证即将展开”。",
        font_size=12,
        color=COLOR_TEXT,
        bold=True,
    )


def main() -> None:
    if not ARCH_IMAGE.exists():
        raise FileNotFoundError(f"architecture image not found: {ARCH_IMAGE}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_cover(prs)
    build_architecture(prs)
    build_innovations(prs)
    build_challenges_and_progress(prs)

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(f"saved {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
