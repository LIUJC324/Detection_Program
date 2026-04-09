from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = PROJECT_ROOT / "docs"
OUTPUT_PPT = DOCS_DIR / "AerialEye_计设比赛答辩汇报_电子科大风格版_20260407.pptx"

IMG_PROJECT_ARCH = DOCS_DIR / "project_architecture.png"
IMG_BACKEND_ARCH = PROJECT_ROOT / "前后端相关" / "后端架构图.png"
IMG_ML_LAYERED = DOCS_DIR / "ml_layered_architecture.png"

SLIDE_W = 13.333
SLIDE_H = 7.5

BLUE = RGBColor(10, 52, 120)
BLUE_DARK = RGBColor(8, 36, 84)
ORANGE = RGBColor(233, 101, 41)
GOLD = RGBColor(245, 176, 65)
GREEN = RGBColor(49, 142, 98)
RED = RGBColor(204, 73, 73)
BG = RGBColor(246, 248, 252)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(218, 224, 235)
TEXT = RGBColor(28, 39, 59)
SUB = RGBColor(89, 103, 126)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(SLIDE_W),
        Inches(SLIDE_H),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    *,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font_name: str = "Microsoft YaHei",
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_top_band(slide, title: str, subtitle: str = "") -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(SLIDE_W),
        Inches(0.78),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    add_text(slide, 0.55, 0.18, 8.5, 0.28, title, 24, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, 9.0, 0.22, 3.7, 0.22, subtitle, 10, color=RGBColor(223, 231, 244), align=PP_ALIGN.RIGHT)


def add_card(slide, left: float, top: float, width: float, height: float, title: str, bullets: list[str], accent: RGBColor) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE
    card.line.width = Pt(1.0)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.14),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, left + 0.28, top + 0.16, width - 0.42, 0.3, title, 18, bold=True)
    body = "\n".join([f"• {item}" for item in bullets])
    add_text(slide, left + 0.28, top + 0.52, width - 0.42, height - 0.65, body, 11, color=SUB)


def add_picture_or_box(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        return
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.color.rgb = LINE
    add_text(slide, left, top + 0.9, width, 0.35, "架构示意", 24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(0.22))
    top.fill.solid()
    top.fill.fore_color.rgb = ORANGE
    top.line.fill.background()

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(0.85),
        Inches(11.9),
        Inches(5.85),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = CARD
    panel.line.color.rgb = LINE

    add_text(slide, 0.98, 1.08, 4.8, 0.24, "电子科技大学风格展示版", 14, color=BLUE, bold=True)
    add_text(slide, 0.98, 1.52, 10.5, 0.62, "AerialEye", 31, color=BLUE_DARK, bold=True)
    add_text(slide, 0.98, 2.12, 11.0, 0.58, "基于 RGB-T 多模态的无人机交通目标检测与联动展示系统", 22, bold=True)
    add_text(slide, 0.98, 2.88, 10.8, 0.42, "从检测模型到平台联动的一体化系统方案", 18, color=SUB)
    add_text(slide, 0.98, 3.46, 8.0, 0.32, "求实求真  大气大为", 16, color=ORANGE, bold=True)

    add_card(
        slide,
        1.0,
        4.18,
        3.45,
        1.62,
        "主题",
        [
            "无人机交通监测",
            "多模态目标检测",
            "系统联动展示",
        ],
        BLUE,
    )
    add_card(
        slide,
        4.75,
        4.18,
        3.45,
        1.62,
        "组成",
        [
            "检测模型",
            "后端平台",
            "流媒体与前端协同",
        ],
        ORANGE,
    )
    add_card(
        slide,
        8.5,
        4.18,
        3.75,
        1.62,
        "特点",
        [
            "模块清晰",
            "链路完整",
            "适合演示与答辩",
        ],
        GOLD,
    )


def build_section(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)

    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(1.35),
        Inches(11.4),
        Inches(4.8),
    )
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor(14, 63, 143)
    block.line.color.rgb = RGBColor(32, 97, 190)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.95),
        Inches(1.35),
        Inches(0.16),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    add_text(slide, 1.35, 2.05, 8.5, 0.65, title, 30, color=WHITE, bold=True)
    add_text(slide, 1.38, 2.92, 8.8, 0.38, subtitle, 16, color=RGBColor(218, 227, 243))


def build_project(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "项目概述")

    add_card(
        slide,
        0.75,
        1.02,
        3.95,
        2.08,
        "场景",
        [
            "无人机视角下的交通目标检测",
            "白天、夜间、低照与复杂天气环境",
            "面向可视化监测和联动展示",
        ],
        BLUE,
    )
    add_card(
        slide,
        4.9,
        1.02,
        3.95,
        2.08,
        "目标",
        [
            "提升复杂环境下的检测稳定性",
            "兼顾小目标表现与系统实时性",
            "支撑前后端联动和固定公网演示",
        ],
        ORANGE,
    )
    add_card(
        slide,
        9.02,
        1.02,
        3.55,
        2.08,
        "成果",
        [
            "模型方案",
            "平台能力",
            "完整演示链路",
        ],
        GOLD,
    )
    add_card(
        slide,
        0.75,
        3.55,
        12.0,
        2.18,
        "系统定义",
        [
            "AerialEye 以 RGB-T 多模态检测为核心，把模型推理、流媒体接入、会话编排、结果回传和前端展示组织为统一系统。",
            "它既是一个检测模型项目，也是一个面向实际应用场景的完整软件系统。",
        ],
        GREEN,
    )


def build_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "系统架构")

    add_picture_or_box(slide, IMG_PROJECT_ARCH, 0.55, 1.0, 7.85, 5.85)
    add_card(
        slide,
        8.7,
        1.12,
        3.95,
        2.0,
        "层次",
        [
            "数据层组织 RGB 与 Thermal 配对输入",
            "模型层完成特征提取、融合和检测输出",
            "服务层提供统一接口与联调入口",
        ],
        BLUE,
    )
    add_card(
        slide,
        8.7,
        3.38,
        3.95,
        2.0,
        "关系",
        [
            "模型关注感知能力",
            "后端关注业务与会话管理",
            "前端负责交互与结果呈现",
        ],
        ORANGE,
    )
    add_text(slide, 8.82, 5.85, 3.6, 0.28, "整体结构按模块协同组织，便于展示与扩展。", 11, color=SUB)


def build_model_scheme(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "模型设计")

    add_picture_or_box(slide, IMG_ML_LAYERED, 0.55, 1.0, 7.2, 5.9)
    add_card(
        slide,
        8.0,
        1.08,
        4.55,
        1.72,
        "双流",
        [
            "RGB 与 Thermal 先分支提特征",
            "避免模态过早混合造成干扰",
        ],
        BLUE,
    )
    add_card(
        slide,
        8.0,
        2.98,
        4.55,
        1.72,
        "融合",
        [
            "跨模态注意力强调互补信息",
            "在夜间和低照场景下提升稳定性",
        ],
        ORANGE,
    )
    add_card(
        slide,
        8.0,
        4.88,
        4.55,
        1.72,
        "增强",
        [
            "BiFPN 做多尺度语义融合",
            "小目标精炼头补充局部细节",
            "FCOS 完成最终检测输出",
        ],
        GOLD,
    )


def build_model_details(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "模型能力")

    add_card(
        slide,
        0.75,
        1.05,
        3.9,
        2.0,
        "输入组织",
        [
            "支持单帧图像对输入",
            "支持拼接视频帧拆分后输入",
            "统一整理为模型推理张量",
        ],
        BLUE,
    )
    add_card(
        slide,
        4.88,
        1.05,
        3.9,
        2.0,
        "输出形式",
        [
            "单帧接口输出标准检测结果",
            "会话接口输出流式推理结果",
            "统一类别映射和模型版本信息",
        ],
        ORANGE,
    )
    add_card(
        slide,
        9.0,
        1.05,
        3.6,
        2.0,
        "特点",
        [
            "面向小目标",
            "适应复杂光照",
            "便于服务化部署",
        ],
        GOLD,
    )
    add_card(
        slide,
        0.75,
        3.45,
        5.9,
        2.18,
        "设计收益",
        [
            "双模态输入提升复杂环境下的鲁棒性。",
            "多尺度结构提升无人机视角小目标的可检测性。",
            "统一推理出口便于和后端系统直接衔接。",
        ],
        GREEN,
    )
    add_card(
        slide,
        6.95,
        3.45,
        5.65,
        2.18,
        "服务适配",
        [
            "模型接口支持健康检查、单帧检测、流会话和视频会话。",
            "既能独立调试，也能作为整套系统中的推理服务使用。",
        ],
        BLUE,
    )


def build_backend(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "后端设计")

    add_picture_or_box(slide, IMG_BACKEND_ARCH, 0.5, 1.0, 7.25, 5.9)
    add_card(
        slide,
        8.0,
        1.08,
        4.55,
        1.68,
        "结构",
        [
            "采用模块化单体组织方式",
            "围绕认证、设备、流媒体、推理、文件和统计拆分职责",
        ],
        BLUE,
    )
    add_card(
        slide,
        8.0,
        2.95,
        4.55,
        2.05,
        "能力",
        [
            "账号与权限管理",
            "设备与双通道绑定",
            "流会话编排与结果推送",
            "统计查询与数据沉淀",
        ],
        ORANGE,
    )
    add_card(
        slide,
        8.0,
        5.22,
        4.55,
        1.25,
        "定位",
        [
            "后端负责把模型能力组织成可管理、可展示、可联动的业务系统。",
        ],
        GREEN,
    )


def build_backend_modules(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "模块划分")

    add_card(slide, 0.72, 1.02, 2.35, 1.9, "认证", ["登录注册", "RBAC 权限", "风控与审计"], BLUE)
    add_card(slide, 3.25, 1.02, 2.35, 1.9, "设备", ["双通道建模", "绑定流程", "在线状态"], ORANGE)
    add_card(slide, 5.78, 1.02, 2.35, 1.9, "流媒体", ["SRS 回调", "会话管理", "WebSocket 推送"], GOLD)
    add_card(slide, 8.3, 1.02, 2.35, 1.9, "推理", ["模型启动", "结果回调", "状态存储"], GREEN)
    add_card(slide, 10.82, 1.02, 1.8, 1.9, "统计", ["明细", "汇总", "趋势"], RED)

    add_card(
        slide,
        0.72,
        3.4,
        12.0,
        2.15,
        "组织方式",
        [
            "模块边界清晰，便于把算法能力嵌入系统。",
            "核心状态通过 Redis、会话与异步任务组织，适合实时展示场景。",
            "业务模块围绕同一条“设备 -> 流 -> 推理 -> 结果 -> 展示”的主链路协同运行。",
        ],
        BLUE,
    )


def build_runtime(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "运行链路")

    add_card(
        slide,
        0.75,
        1.05,
        3.9,
        2.05,
        "采集",
        [
            "摄像头或视频源进入平台",
            "RGB 与 IR 按规则组织",
            "支持实时流和离线视频",
        ],
        BLUE,
    )
    add_card(
        slide,
        4.88,
        1.05,
        3.9,
        2.05,
        "推理",
        [
            "模型端主动拉流或读取视频",
            "拆分模态、完成检测、生成结果",
            "提供单帧与会话两类输出",
        ],
        ORANGE,
    )
    add_card(
        slide,
        9.0,
        1.05,
        3.6,
        2.05,
        "展示",
        [
            "后端接收结果并分发",
            "前端完成画面叠加",
            "形成完整演示闭环",
        ],
        GOLD,
    )
    add_card(
        slide,
        0.75,
        3.45,
        12.0,
        2.18,
        "部署方式",
        [
            "模型服务运行在本机，后端与流媒体运行在云端，固定公网入口统一承接模型调用。",
            "这种形态兼顾模型计算能力与现场演示稳定性，也便于后续把系统拆分到更大规模的部署环境。",
        ],
        GREEN,
    )
    add_text(slide, 0.85, 6.05, 11.8, 0.3, "固定公网入口可直接支撑外部联调和现场展示。", 12, color=SUB)


def build_showcase(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "展示模块")

    add_card(
        slide,
        0.75,
        1.02,
        3.95,
        2.0,
        "模型接口",
        [
            "健康检查",
            "模型信息查询",
            "单帧检测与会话控制",
        ],
        BLUE,
    )
    add_card(
        slide,
        4.88,
        1.02,
        3.95,
        2.0,
        "实时展示",
        [
            "视频画面叠加结果",
            "状态与结果同步更新",
            "面向演示的交互组织",
        ],
        ORANGE,
    )
    add_card(
        slide,
        9.0,
        1.02,
        3.6,
        2.0,
        "管理视图",
        [
            "设备、会话、统计",
            "记录与回看",
            "系统状态查询",
        ],
        GOLD,
    )
    add_card(
        slide,
        0.75,
        3.42,
        12.0,
        2.2,
        "展示重点",
        [
            "这套系统适合按“项目场景 -> 系统架构 -> 模型方案 -> 平台能力 -> 演示链路”的顺序进行展示。",
            "这样既能讲清楚技术设计，也能把完整作品的系统价值呈现出来。",
        ],
        GREEN,
    )


def build_value(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    add_top_band(slide, "总结与展望")

    add_card(
        slide,
        0.75,
        1.08,
        5.9,
        4.8,
        "项目价值",
        [
            "把多模态检测算法与完整系统能力结合到同一作品中。",
            "不仅能输出检测结果，还能围绕设备、流、会话、回调和展示形成闭环。",
            "更接近真实应用系统，而不是只停留在单个模型页面。",
        ],
        GREEN,
    )
    add_card(
        slide,
        6.9,
        1.08,
        5.7,
        4.8,
        "扩展方向",
        [
            "继续提升多场景下的识别稳定性与推理效率。",
            "完善前端展示层与统计分析能力。",
            "增强 GPU 加速、并发支持和更复杂的部署能力。",
        ],
        ORANGE,
    )
    add_text(
        slide,
        0.9,
        6.22,
        11.6,
        0.32,
        "AerialEye 的核心价值，在于把模型能力、平台能力和展示能力组织成一个完整作品。",
        14,
        color=BLUE_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.1),
        Inches(1.28),
        Inches(11.1),
        Inches(4.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(14, 63, 143)
    panel.line.color.rgb = RGBColor(31, 97, 191)

    add_text(slide, 1.58, 2.08, 4.0, 0.45, "结束", 28, color=WHITE, bold=True)
    add_text(slide, 1.58, 3.0, 5.2, 0.52, "AerialEye", 31, color=RGBColor(255, 215, 175), bold=True)
    add_text(slide, 1.58, 3.82, 9.5, 0.55, "多模态检测、后端工程与联动展示组成的完整系统方案", 18, color=RGBColor(220, 228, 244))
    add_text(slide, 1.58, 5.12, 3.0, 0.28, "谢谢", 18, color=ORANGE, bold=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_cover(prs)
    build_section(prs, "项目", "场景、目标与系统定义")
    build_project(prs)
    build_section(prs, "架构", "整体结构与模块关系")
    build_architecture(prs)
    build_section(prs, "模型", "方案设计与能力组织")
    build_model_scheme(prs)
    build_model_details(prs)
    build_section(prs, "后端", "平台能力与工程设计")
    build_backend(prs)
    build_backend_modules(prs)
    build_section(prs, "运行", "部署方式与展示链路")
    build_runtime(prs)
    build_showcase(prs)
    build_section(prs, "总结", "项目价值与扩展方向")
    build_value(prs)
    build_closing(prs)

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(f"saved {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
